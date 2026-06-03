from pptx import Presentation
from pptx.util import Inches


def create_pitch_deck(report: dict, filename="pitch_deck.pptx"):

    prs = Presentation()

    def add_slide(title, content):
        slide_layout = prs.slides.add_slide(prs.slide_layouts[1])
        slide = prs.slides[-1]

        slide.shapes.title.text = title
        slide.placeholders[1].text = str(content)

    # 1. Problem
    add_slide("Problem", report.get("market", {}).get("problem", "Not available"))

    # 2. Solution
    add_slide("Solution", report.get("market", {}).get("solution", "Not available"))

    # 3. Market
    add_slide("Market Analysis", report.get("market", "N/A"))

    # 4. Competitors
    add_slide("Competitors", report.get("competitors", "N/A"))

    # 5. SWOT
    add_slide("SWOT Analysis", report.get("swot", "N/A"))

    # 6. Revenue Model
    add_slide("Revenue Model", report.get("revenue", "N/A"))

    # 7. Investor Readiness
    add_slide("Investor Score", report.get("investor", "N/A"))

    prs.save(filename)

    return filename